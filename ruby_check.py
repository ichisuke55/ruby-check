import re
import pptx
import json
import argparse


# Regular expression for kanji.
# ref: https://note.nkmk.me/python-re-regex-character-type/
KANJI_REGEXP = \
    ('[\u2E80-\u2FDF\u3005-\u3007\u3400-\u4DBF\u4E00-'
     '\u9FFF\uF900-\uFAFF\U00020000-\U0002EBEF]+')
RED = '\033[31m'
END = '\033[0m'


def load_json_file():
    with open('./kanji.json') as f:
        d = json.load(f)
    return d


def ruby_check(text: str, hisshu_list: list) -> str:
    match_list = re.findall(KANJI_REGEXP, text)
    if len(match_list) > 0:
        eval_text = "".join(match_list)
        for t in eval_text:
            if t not in hisshu_list:
                return text.replace(t, f'{RED}{t}{END}')


def main():
    # Define argument parser
    parser = argparse.ArgumentParser()
    parser.add_argument('-f', '--file',
                        help='specify pptx file path',
                        required=True)
    parser.add_argument('-g', '--grade',
                        help='specify elementary school grades (default: 3)',
                        type=int, choices=range(1, 7), default=3)
    args = parser.parse_args()
    ppt_obj = pptx.Presentation(args.file)

    # Load grade and kanji key-values
    kanji_dict = load_json_file()

    limit_grade = args.grade
    kanji_list = [v for k, v in kanji_dict.items() if int(k) <= limit_grade]
    kanji_hisshu_list = "".join(kanji_list)

    for i, slide in enumerate(ppt_obj.slides, start=1):
        print(f'=== page{i} ===')
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = ruby_check(shape.text, kanji_hisshu_list)
                if text is not None:
                    print(text)


if __name__ == '__main__':
    main()
